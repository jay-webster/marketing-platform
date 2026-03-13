"""Format-specific text extraction utilities for the ingestion pipeline.

All sync extraction functions should be called via asyncio.to_thread() at call sites.
"""
import asyncio
import csv
import io
import zipfile
from typing import Callable

MAX_FILE_SIZE_BYTES = 50 * 1024 * 1024  # 50 MB

# Human-readable failure reasons (FR-6.2)
REASON_EMPTY = "File is empty."
REASON_CORRUPT = "File could not be read — it may be corrupted."
REASON_NO_TEXT = "No readable text content found — the file may contain only images."
REASON_OVERSIZED = "File exceeds the 50 MB size limit."


def preflight_check(size_bytes: int, filename: str) -> None:
    """Raise ValueError for empty or oversized files before extraction."""
    if size_bytes == 0:
        raise ValueError(REASON_EMPTY)
    if size_bytes > MAX_FILE_SIZE_BYTES:
        raise ValueError(REASON_OVERSIZED)


def _extract_pdf(stream: io.BytesIO) -> str:
    try:
        import fitz  # pymupdf  # noqa: PLC0415
    except ImportError as exc:
        raise RuntimeError("pymupdf is not installed") from exc

    try:
        doc = fitz.open(stream=stream.read(), filetype="pdf")
    except Exception as exc:
        raise ValueError(REASON_CORRUPT) from exc

    if doc.page_count == 0:
        raise ValueError(REASON_EMPTY)

    pages: list[str] = []
    all_image_only = True
    for page in doc:
        text = page.get_text("text")
        if text.strip():
            all_image_only = False
            pages.append(text)
        else:
            images = page.get_images()
            if images:
                pages.append(f"[IMAGE-ONLY PAGE {page.number + 1}]")
            else:
                pages.append("")

    doc.close()

    if all_image_only:
        raise ValueError(REASON_NO_TEXT)

    result = "\n\n".join(p for p in pages if p)
    if not result.strip():
        raise ValueError(REASON_NO_TEXT)
    return result


def _extract_docx(stream: io.BytesIO) -> str:
    try:
        from docx import Document  # noqa: PLC0415
        from docx.opc.exceptions import PackageNotFoundError  # noqa: PLC0415
    except ImportError as exc:
        raise RuntimeError("python-docx is not installed") from exc

    try:
        doc = Document(stream)
    except (zipfile.BadZipFile, Exception) as exc:
        raise ValueError(REASON_CORRUPT) from exc

    lines: list[str] = []

    for para in doc.paragraphs:
        style = para.style.name if para.style else ""
        text = para.text.strip()
        if not text:
            continue
        if style.startswith("Heading 1"):
            lines.append(f"# {text}")
        elif style.startswith("Heading 2"):
            lines.append(f"## {text}")
        elif style.startswith("Heading 3"):
            lines.append(f"### {text}")
        elif "List" in style:
            lines.append(f"- {text}")
        else:
            lines.append(text)

    for table in doc.tables:
        for row in table.rows:
            lines.append(" | ".join(cell.text.strip() for cell in row.cells))

    if not lines:
        raise ValueError(REASON_EMPTY)
    return "\n".join(lines)


def _extract_pptx(stream: io.BytesIO) -> str:
    try:
        from pptx import Presentation  # noqa: PLC0415
    except ImportError as exc:
        raise RuntimeError("python-pptx is not installed") from exc

    try:
        prs = Presentation(stream)
    except (zipfile.BadZipFile, Exception) as exc:
        raise ValueError(REASON_CORRUPT) from exc

    slides_text: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        parts = [f"## Slide {i}"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        parts.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    parts.append(" | ".join(cell.text.strip() for cell in row.cells))
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append(f"[Notes: {notes}]")
        slides_text.append("\n".join(parts))

    if not slides_text:
        raise ValueError(REASON_EMPTY)
    return "\n\n".join(slides_text)


def _extract_csv(stream: io.BytesIO) -> str:
    try:
        content = stream.read().decode("utf-8", errors="replace")
    except Exception as exc:
        raise ValueError(REASON_CORRUPT) from exc

    if not content.strip():
        raise ValueError(REASON_EMPTY)

    try:
        reader = csv.reader(io.StringIO(content))
        rows = list(reader)
    except csv.Error as exc:
        raise ValueError(REASON_CORRUPT) from exc

    if not rows:
        raise ValueError(REASON_EMPTY)

    lines: list[str] = []
    header = rows[0]
    lines.append(" | ".join(header))
    lines.append(" | ".join(["---"] * len(header)))
    for row in rows[1:]:
        lines.append(" | ".join(row))
    return "\n".join(lines)


def _extract_text(stream: io.BytesIO) -> str:
    try:
        content = stream.read().decode("utf-8", errors="replace")
    except Exception as exc:
        raise ValueError(REASON_CORRUPT) from exc

    if not content.strip():
        raise ValueError(REASON_EMPTY)
    return content


_EXTRACTORS: dict[str, Callable[[io.BytesIO], str]] = {
    ".pdf": _extract_pdf,
    ".docx": _extract_docx,
    ".pptx": _extract_pptx,
    ".csv": _extract_csv,
    ".txt": _extract_text,
    ".md": _extract_text,
}

SUPPORTED_EXTENSIONS = set(_EXTRACTORS.keys())


def extract_text(stream: io.BytesIO, file_type: str) -> str:
    """Dispatch extraction by file extension. Synchronous — wrap with asyncio.to_thread()."""
    ext = file_type.lower()
    handler = _EXTRACTORS.get(ext)
    if handler is None:
        raise ValueError(f"Unsupported file type: {ext}")
    return handler(stream)


async def extract_text_async(stream: io.BytesIO, file_type: str) -> str:
    """Async wrapper for extract_text."""
    return await asyncio.to_thread(extract_text, stream, file_type)
